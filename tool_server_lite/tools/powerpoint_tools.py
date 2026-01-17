#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PowerPoint 操作工具
"""

from pathlib import Path
from typing import Dict, Any, Optional
from .file_tools import BaseTool, get_abs_path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from PIL import Image
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class ImagesToPptTool(BaseTool):
    """将多张图片转换为 PPT - 每张图片一页，PPT 尺寸自动匹配第一张图片"""
    
    def execute(self, task_id: str, parameters: Dict[str, Any]) -> Dict[str, Any]:
        """
        将多张图片转换为 PowerPoint（简化版）
        
        Parameters:
            image_paths (list[str]): 图片文件相对路径列表
            output_path (str): 输出 PPT 文件相对路径
        
        Returns:
            status: "success" 或 "error"
            output: PPT 文件保存位置信息
            error: 错误信息（如有）
        
        Note:
            - PPT 尺寸会自动根据第一张图片的宽高比设置
            - 每张图片占满整页（保持比例，居中显示）
        """
        try:
            if not PPTX_AVAILABLE:
                return {
                    "status": "error",
                    "output": "",
                    "error": "python-pptx 未安装。请运行: pip install python-pptx"
                }
            
            # 获取参数
            image_paths = parameters.get("image_paths")
            output_path = parameters.get("output_path")
            
            if not image_paths:
                return {
                    "status": "error",
                    "output": "",
                    "error": "缺少必需参数: image_paths"
                }
            
            if not output_path:
                return {
                    "status": "error",
                    "output": "",
                    "error": "缺少必需参数: output_path"
                }
            
            # 统一转换为列表
            if isinstance(image_paths, str):
                image_paths = [image_paths]
            
            # 转换为绝对路径
            abs_image_paths = [get_abs_path(task_id, path) for path in image_paths]
            abs_output_path = get_abs_path(task_id, output_path)
            
            # 检查图片是否存在
            for img_path in abs_image_paths:
                if not img_path.exists():
                    return {
                        "status": "error",
                        "output": "",
                        "error": f"图片文件不存在: {img_path.name}"
                    }
            
            # 确保输出目录存在
            abs_output_path.parent.mkdir(parents=True, exist_ok=True)
            
            # 读取第一张图片，确定 PPT 尺寸
            first_img = Image.open(abs_image_paths[0])
            img_width, img_height = first_img.size
            img_ratio = img_width / img_height
            
            # 根据图片比例设置 PPT 尺寸
            # 标准宽度 10 英寸，高度根据比例计算
            slide_width = 10.0
            slide_height = slide_width / img_ratio
            
            print(f"[INFO] PPT 尺寸: {slide_width:.2f}\" x {slide_height:.2f}\" (根据图片 {img_width}x{img_height} 计算)")
            
            # 创建 PPT
            prs = Presentation()
            prs.slide_width = Inches(slide_width)
            prs.slide_height = Inches(slide_height)
            
            # 添加图片页
            for idx, img_path in enumerate(abs_image_paths, 1):
                # 使用空白布局
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                
                # 获取当前图片尺寸
                img = Image.open(img_path)
                img_w, img_h = img.size
                current_ratio = img_w / img_h
                ppt_ratio = slide_width / slide_height
                
                # 计算图片位置和尺寸（保持比例，适应页面）
                if abs(current_ratio - ppt_ratio) < 0.01:
                    # 比例相同，填满整页
                    left = Inches(0)
                    top = Inches(0)
                    width = Inches(slide_width)
                    height = Inches(slide_height)
                else:
                    # 比例不同，保持比例居中
                    if current_ratio > ppt_ratio:
                        # 图片更宽，以宽度为准
                        width = Inches(slide_width)
                        height = Inches(slide_width / current_ratio)
                        left = Inches(0)
                        top = Inches((slide_height - height.inches) / 2)
                    else:
                        # 图片更高，以高度为准
                        height = Inches(slide_height)
                        width = Inches(slide_height * current_ratio)
                        top = Inches(0)
                        left = Inches((slide_width - width.inches) / 2)
                
                # 添加图片到幻灯片
                slide.shapes.add_picture(
                    str(img_path),
                    left, top,
                    width=width,
                    height=height
                )
                
                print(f"[INFO] 添加图片 {idx}/{len(abs_image_paths)}: {img_path.name}")
            
            # 保存 PPT
            prs.save(str(abs_output_path))
            
            file_size = abs_output_path.stat().st_size / 1024  # KB
            
            return {
                "status": "success",
                "output": f"PPT 已生成: {output_path}\n- 总页数: {len(prs.slides)}\n- 图片数: {len(abs_image_paths)}\n- 文件大小: {file_size:.1f} KB",
                "error": ""
            }
            
        except Exception as e:
            return {
                "status": "error",
                "output": "",
                "error": f"生成 PPT 失败: {str(e)}"
            }
